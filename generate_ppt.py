from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

def add_slide(prs, title_text, bullet_points):
    """Utility to add a standard slide with title and bullets."""
    # 1 is the layout for Title and Content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    title = slide.shapes.title
    title.text = title_text
    
    # Set Body
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()  # Clear existing default text
    
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph()
        p.text = point
        # Add a bit of space after paragraphs
        p.space_after = Pt(14)
        p.font.size = Pt(20)
        p.level = 0
        if ":" in point:     # Make prefix bold if formatted like "Definition: ..."
            prefix = point.split(':')[0] + ":"
            rest = point[len(prefix):]
            p.clear()
            run1 = p.add_run()
            run1.text = prefix
            run1.font.bold = True
            run2 = p.add_run()
            run2.text = rest
            
    return slide

# =========================================================
# Slide 1: Title Slide
# =========================================================
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Stock Trend Prediction Using Python & Machine Learning (LSTM)\nwith Flask Web Application"
title.text_frame.paragraphs[0].font.size = Pt(36)
subtitle.text = "Student Name: [Your Name]\nCourse: [Your Course / Degree]\nCollege: [Your College / University Name]"

# =========================================================
# Content Slides
# =========================================================
slides_data = [
    {
        "title": "What is the Stock Market?",
        "bullets": [
            "Definition: A decentralized market for buying and selling publicly held company shares.",
            "Importance: A primary source for companies to raise capital.",
            "Economic Indicator: Serves as a barometer for the broader economy's health.",
            "Volatility: Prices fluctuate based on earnings, news, sentiment, and global events."
        ]
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "High Volatility: Stock prices change rapidly by the second due to thousands of unquantifiable external factors.",
            "Non-Linear Data: Financial time-series data is highly complex, noisy, and non-linear.",
            "Traditional Flaws: Basic statistical models are ineffective against modern market speeds.",
            "Emotional Trading: Human sentiment and news cycles heavily dictate market movements.",
            "The Need: Investors need an analytical, emotionless system to extract mathematical trends."
        ]
    },
    {
        "title": "Project Objectives",
        "bullets": [
            "Accurate Forecasting: Build an advanced LSTM model to predict future prices based on historical trends.",
            "Seamless Accessibility: Develop a user-friendly Web Dashboard for global/Indian stock analysis.",
            "Real-time Insights: Implement live tracking of market movers (NIFTY 50 / SENSEX / SPY).",
            "Automated Signals: Output actionable financial heuristics like 'BUY' or 'SELL'."
        ]
    },
    {
        "title": "Overview of Stock Trend Prediction",
        "bullets": [
            "Traditional vs. ML Analysis: ML discovers hidden correlation patterns autonomously.",
            "Time-Series Forecasting: Predicting future values based on previously observed, sequential values.",
            "Feature Extraction: ML models learn to weigh features such as 'Open', 'High', 'Low', 'Close', and 'Volume'.",
            "Continuous Learning: The system adapts to recent historical volatility for near-term predictions."
        ]
    },
    {
        "title": "Technologies Used",
        "bullets": [
            "Programming Language: Python 3 (Robust scientific and Web libraries).",
            "Deep Learning Framework: TensorFlow & Keras (Building the Neural Network).",
            "Backend Web Framework: Flask (Server logic) & SQLite (Database mapping).",
            "Data API: yfinance (Direct feed from Yahoo Finance for live pricing).",
            "Frontend UI: HTML5, Bootstrap 5 UI, Chart.js & Plotly for interactive data viz."
        ]
    },
    {
        "title": "System Architecture",
        "bullets": [
            "User Input (Frontend): The user interactively selects a stock on the dashboard.",
            "API Fetch (Backend): Flask calls the API to retrieve 2 years of daily financial numbers.",
            "Data Processing: Python normalizes the data into fixed sequences.",
            "Prediction Engine: The data flows into the ML algorithm to calculate predictive outcomes.",
            "Render (Frontend): Plotly maps the mathematical arrays back into interactive web charts."
        ]
    },
    {
        "title": "Data Collection",
        "bullets": [
            "Primary Source: Yahoo Finance & Internal NLP Engines.",
            "Scope: 1 to 2 years of historical daily data per query.",
            "Attributes Collected: Close Price, Volume, Open/High/Low ranges.",
            "Advantage: Live on-demand fetching prevents stale model architectures."
        ]
    },
    {
        "title": "Data Preprocessing",
        "bullets": [
            "Handling Missing Data: Using backfill methods to seamlessly patch empty market days.",
            "Feature Scaling: Converting stock prices to a scale between 0 and 1 using MinMaxScaler.",
            "Sequencing Data: Creating chronological sequences (e.g., 60 days) to predict the 61st day.",
            "Train/Test Split: 80% of historical data for training, 20% held back to test accuracy."
        ]
    },
    {
        "title": "Introduction to LSTM",
        "bullets": [
            "Definition: Long Short-Term Memory is an advanced Recurrent Neural Network (RNN).",
            "The Problem with Standard RNNs: They suffer from the 'Vanishing Gradient Problem' and 'forget' older data.",
            "Why LSTM?: Capable of storing relevant seqeuence information in 'Memory Cells'.",
            "Mechanism: Ideal for discarding irrelevant daily noise while keeping long-term financial structures."
        ]
    },
    {
        "title": "Working of LSTM",
        "bullets": [
            "Cell State: The 'conveyor belt' holding long-term memory across sequence steps.",
            "Forget Gate: Decides what information from the previous frame should be discarded.",
            "Input Gate: Decides what new data to store in the cell state.",
            "Output Gate: Decides the next hidden state based on the updated cell memory.",
            "Result: Dynamically identifies patterns without human bias."
        ]
    },
    {
        "title": "Model Training",
        "bullets": [
            "Architecture: Stacking sequential layers of Long Short-Term Memory units.",
            "Activation Function: Rectified Linear Unit (ReLU) to handle complex, non-linear patterns.",
            "Optimization: The 'Adam' optimizer tweaks network weights incrementally.",
            "Loss Metric: Mean Squared Error (MSE) measures prediction errors.",
            "Epochs & Batching: Data is iteratively passed back to reduce total variance."
        ]
    },
    {
        "title": "Prediction Results",
        "bullets": [
            "Test Accuracy Alignment: The model tests itself on the 20% unseen data block.",
            "Visual Verification: Actual Price vs Predicted Price are mapped dynamically.",
            "Key Metric: Root Mean Square Error (RMSE) acts as the mathematical exactness score.",
            "Future Forecast: Projecting the last 60 days of data autonomously up to 30 unseen days ahead."
        ]
    },
    {
        "title": "Indian Stocks Listing Feature",
        "bullets": [
            "Regional Support: Automated sweeping of NIFTY 50 and SENSEX equities.",
            "Dynamic Parsing: Live calculation of Daily Gain/Loss percentages.",
            "User Access: Features rapid search filtering for thousands of companies.",
            "Actionable UI: 1-Click routing from the Market Tracker directly to the Deep Learning engine."
        ]
    },
    {
        "title": "Dashboard UI",
        "bullets": [
            "Aesthetic Design: Built on Bootstrap 5 with glassmorphism overlays and custom 'Outfit' typography.",
            "Responsiveness: Flawless layout across Desktop, Tablet, and Mobile views.",
            "Theme Switching: Real-time Light/Dark mode toggling affecting CSS and JS Charts.",
            "Visual Elements: Hover animations, dynamic shadows, and gradient heading fonts."
        ]
    },
    {
        "title": "Advanced Features",
        "bullets": [
            "Technical Overlays: Auto-generation of 20-Day Simple (MA20) & 50-Day Exponential (EMA50) Moving Averages.",
            "News Sentiment Analysis: NLP scanning of Yahoo Finance headlines to grade sentiment (Positive/Neutral/Negative).",
            "Account Persistence: Secure SQLite database schema securely storing hashed user passwords.",
            "Custom Dashboard: Saving favorite tracked stocks to a personal 'Watchlist'."
        ]
    },
    {
        "title": "Advantages of the System",
        "bullets": [
            "Time Efficiency: Compresses hours of chart analysis into milliseconds.",
            "Emotionless Execution: Strictly objective AI-driven logic protecting from fear and greed.",
            "High Scope Integration: Scales securely across international (.US, .NS, .UK) stock symbols.",
            "Cost Effectiveness: Enterprise-grade insights built on open-source Python frameworks."
        ]
    },
    {
        "title": "Limitations",
        "bullets": [
            "'Black Swan' Events: AI relies on historical precedent and cannot predict sudden pandemics or CEO scandals.",
            "Lagging Fundamentals: The network tracks price action, completely blind to P/E ratios or corporate debt.",
            "Not Financial Advice: Unprecedented real-world liquidity crunches can quickly violate mathematical limits."
        ]
    },
    {
        "title": "Future Scope",
        "bullets": [
            "Fundamental Integration: Feeding the model quarterly Balance Sheet data alongside raw price history.",
            "Push Notifications: Integrating SMS/Email alerts for massive automated 'BUY' signal triggers.",
            "Cloud Deployment: Migrating the SQLite core to PostgreSQL managed by AWS/Docker instances.",
            "Broker Hooks: Real-time API integration with brokers (Zerodha/Upstox) for automated order executions."
        ]
    },
    {
        "title": "Conclusion",
        "bullets": [
            "We successfully merged classical Deep Learning algorithms with a modern Web Dashboard.",
            "The LSTM module actively discovers non-evident patterns in highly volatile market frames.",
            "Empowers retail investors to utilize logic previously restricted to institutional, AI-driven strategies."
        ]
    }
]

for slide_data in slides_data:
    add_slide(prs, slide_data["title"], slide_data["bullets"])

# =========================================================
# Final Slide: Questions
# =========================================================
final_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(final_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Thank You!"
subtitle.text = "Any Questions?"

# Save presentation
prs.save('Stock_Trend_Prediction_Presentation.pptx')
print("Presentation generated successfully at Stock_Trend_Prediction_Presentation.pptx")
